"""
PPT Engine — renders SlideDeck JSON to editable .pptx files via python-pptx.
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

from agent.ppt_engine.dsl_schema import SlideDeck, Slide, ChartData


# Theme color palettes
THEMES = {
    "academic_blue": {
        "title_bg": RGBColor(0x00, 0x3C, 0x71),
        "title_fg": RGBColor(0xFF, 0xFF, 0xFF),
        "body_fg": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0x00, 0x7B, 0xC0),
    },
    "business_gray": {
        "title_bg": RGBColor(0x2D, 0x2D, 0x2D),
        "title_fg": RGBColor(0xFF, 0xFF, 0xFF),
        "body_fg": RGBColor(0x33, 0x33, 0x33),
        "accent": RGBColor(0xE8, 0x4E, 0x0E),
    },
}

CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}


def render_pptx(deck: SlideDeck, output_path: str) -> str:
    """Render a SlideDeck to a .pptx file. Returns the output file path."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    theme = THEMES.get(deck.meta.theme, THEMES["academic_blue"])

    for slide_data in deck.slides:
        _add_slide(prs, slide_data, theme)

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return output_path


def _add_slide(prs: Presentation, s: Slide, theme: dict):
    """Add a single slide based on layout type."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    if s.layout == "title_slide":
        _render_title_slide(slide, s, theme)
    elif s.layout == "section_header":
        _render_section_header(slide, s, theme)
    elif s.layout == "content_with_chart":
        _render_content_with_chart(slide, s, theme)
    elif s.layout == "two_column":
        _render_two_column(slide, s, theme)
    elif s.layout == "comparison":
        _render_two_column(slide, s, theme)  # same layout, different label
    else:
        # content_only, content_with_image, summary — all text-based
        _render_content_only(slide, s, theme)

    # Speaker notes
    if s.speaker_notes:
        slide.notes_slide.notes_text_frame.text = s.speaker_notes


def _render_title_slide(slide, s: Slide, theme: dict):
    # Full background color
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = theme["title_bg"]

    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s.title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = theme["title_fg"]
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if s.subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11), Inches(1))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = s.subtitle
        p2.font.size = Pt(24)
        p2.font.color.rgb = theme["title_fg"]
        p2.alignment = PP_ALIGN.CENTER


def _render_section_header(slide, s: Slide, theme: dict):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = theme["accent"]

    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = s.title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = theme["title_fg"]
    p.alignment = PP_ALIGN.CENTER


def _render_content_only(slide, s: Slide, theme: dict):
    # Title bar
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = s.title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["accent"]

    # Body
    if s.body:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(s.body.split("\n")):
            if i == 0:
                tf2.paragraphs[0].text = line
                tf2.paragraphs[0].font.size = Pt(18)
                tf2.paragraphs[0].font.color.rgb = theme["body_fg"]
            else:
                p = tf2.add_paragraph()
                p.text = line
                p.font.size = Pt(18)
                p.font.color.rgb = theme["body_fg"]

    # Image prompt placeholder
    if s.image_prompt:
        txBox3 = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(4.5), Inches(4))
        tf3 = txBox3.text_frame
        tf3.word_wrap = True
        tf3.paragraphs[0].text = f"[Image: {s.image_prompt}]"
        tf3.paragraphs[0].font.size = Pt(12)
        tf3.paragraphs[0].font.italic = True
        tf3.paragraphs[0].font.color.rgb = RGBColor(0x99, 0x99, 0x99)


def _render_content_with_chart(slide, s: Slide, theme: dict):
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = s.title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["accent"]

    # Body text (left side)
    if s.body:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.5), Inches(5.5))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for i, line in enumerate(s.body.split("\n")):
            if i == 0:
                tf2.paragraphs[0].text = line
                tf2.paragraphs[0].font.size = Pt(16)
                tf2.paragraphs[0].font.color.rgb = theme["body_fg"]
            else:
                p2 = tf2.add_paragraph()
                p2.text = line
                p2.font.size = Pt(16)
                p2.font.color.rgb = theme["body_fg"]

    # Chart (right side)
    if s.chart:
        _add_chart(slide, s.chart, Inches(6.5), Inches(1.5), Inches(6), Inches(5))


def _render_two_column(slide, s: Slide, theme: dict):
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = s.title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = theme["accent"]

    # Left column
    left_text = s.body_left or s.body or ""
    txBox_l = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    for i, line in enumerate(left_text.split("\n")):
        if i == 0:
            tf_l.paragraphs[0].text = line
            tf_l.paragraphs[0].font.size = Pt(16)
        else:
            p = tf_l.add_paragraph()
            p.text = line
            p.font.size = Pt(16)

    # Right column
    right_text = s.body_right or ""
    txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    for i, line in enumerate(right_text.split("\n")):
        if i == 0:
            tf_r.paragraphs[0].text = line
            tf_r.paragraphs[0].font.size = Pt(16)
        else:
            p = tf_r.add_paragraph()
            p.text = line
            p.font.size = Pt(16)


def _add_chart(slide, chart: ChartData, left, top, width, height):
    """Add a chart shape to the slide."""
    chart_type = CHART_TYPE_MAP.get(chart.type, XL_CHART_TYPE.COLUMN_CLUSTERED)
    cd = CategoryChartData()

    labels = chart.data.get("labels", [])
    cd.categories = labels

    # Support single series (values) or multi-series (series)
    if "values" in chart.data:
        series_name = chart.unit or "Value"
        cd.add_series(series_name, chart.data["values"])
    elif "series" in chart.data:
        for s in chart.data["series"]:
            cd.add_series(s["name"], s["values"])

    slide.shapes.add_chart(chart_type, left, top, width, height, cd)
